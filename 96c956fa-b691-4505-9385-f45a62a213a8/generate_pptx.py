from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Cm(29.7)
prs.slide_height = Cm(21.0)
blank_slide_layout = prs.slide_layouts[6]

dark_blue = RGBColor(10, 30, 63)
cyan = RGBColor(0, 155, 214)
white = RGBColor(255, 255, 255)

def add_text(slide, text, left, top, width, height, font_size=Pt(12), bold=False, color=RGBColor(0,0,0), align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox

# SLIDE 1: EXTERIOR
slide1 = prs.slides.add_slide(blank_slide_layout)

# Panel 1 (Left): Flap
add_text(slide1, "NUESTRA PROPUESTA", Cm(1), Cm(2), Cm(7.9), Cm(1), font_size=Pt(16), bold=True, color=cyan)
add_text(slide1, "El sistema educativo dominicano genera millones de datos a través del Registro de Grado. SIGERA digitaliza este proceso para convertir el trabajo manual en información estratégica, garantizando precisión, transparencia y agilidad.", Cm(1), Cm(3.5), Cm(7.9), Cm(4), font_size=Pt(10))
add_text(slide1, "• Cálculo exacto de promedios de competencias.\n• Automatización de Completivas y Extraordinarias.\n• Alertas preventivas de inasistencia.\n• Boletines con formato oficial MINERD.\n• Cumplimiento de la Ordenanza 04-2023.", Cm(1), Cm(7.5), Cm(7.9), Cm(4), font_size=Pt(10))
try:
    slide1.shapes.add_picture("scratch/pptx_extracted/ppt/media/image2.png", Cm(1), Cm(13), width=Cm(7.9))
except Exception as e:
    pass

# Panel 2 (Center): Back Cover
shape = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(9.9), Cm(0), Cm(9.9), Cm(11))
shape.fill.solid()
shape.fill.fore_color.rgb = dark_blue
shape.line.fill.background()

add_text(slide1, "EL PODER DE LOS DATOS\nAL SERVICIO DE\nLA EDUCACIÓN", Cm(10.5), Cm(3), Cm(8.9), Cm(4), font_size=Pt(18), bold=True, color=white, align=PP_ALIGN.CENTER)

add_text(slide1, "CONTACTO", Cm(10.5), Cm(12), Cm(8.9), Cm(1), font_size=Pt(16), bold=True, color=cyan)
add_text(slide1, "SEDE CENTRAL\nAv. Winston Churchill #123\n\nTELÉFONOS\nOficina: +1 809-555-0199\n\nEN LÍNEA\ninfo@sigera.edu.do", Cm(10.5), Cm(13.5), Cm(8.9), Cm(6), font_size=Pt(10))


# Panel 3 (Right): Front Cover
try:
    slide1.shapes.add_picture("scratch/pptx_extracted/ppt/media/image31.jpeg", Cm(19.8), Cm(0), width=Cm(9.9), height=Cm(12))
except Exception as e:
    pass

add_text(slide1, "SIGERA", Cm(20.5), Cm(13), Cm(8.5), Cm(2), font_size=Pt(36), bold=True, color=cyan)
add_text(slide1, "SISTEMA DE GESTIÓN", Cm(20.5), Cm(15), Cm(8.5), Cm(1), font_size=Pt(18), bold=True, color=dark_blue)
add_text(slide1, "Solución integral para el registro y análisis del rendimiento académico en el nivel secundario. Cumplimiento absoluto con los estándares del MINERD.", Cm(20.5), Cm(16.5), Cm(8.5), Cm(3), font_size=Pt(10))

# SLIDE 2: INTERIOR
slide2 = prs.slides.add_slide(blank_slide_layout)

# Panel 4 (Left)
add_text(slide2, "OBJETIVOS CLAVE", Cm(1), Cm(8), Cm(7.9), Cm(1), font_size=Pt(16), bold=True, color=cyan)
add_text(slide2, "SIGERA fue diseñado para resolver problemáticas reales de la gestión académica dominicana.\n\nEficiencia Docente:\nEliminación del trabajo repetitivo en cálculos y transcripciones hacia los boletines de calificaciones.\n\nDecisiones Basadas en Datos:\nLa dirección escolar cuenta con analítica detallada de aprobación y riesgo escolar.", Cm(1), Cm(9.5), Cm(7.9), Cm(8), font_size=Pt(10))

# Panel 5 (Center)
add_text(slide2, "NUEVAS POSIBILIDADES", Cm(10.5), Cm(2), Cm(8.9), Cm(1), font_size=Pt(16), bold=True, color=cyan)
add_text(slide2, "El sistema interconecta a toda la comunidad educativa. Los coordinadores pueden validar registros en tiempo real, mientras los directores visualizan el progreso del centro.", Cm(10.5), Cm(3.5), Cm(8.9), Cm(3), font_size=Pt(10))

try:
    slide2.shapes.add_picture("scratch/pptx_extracted/ppt/media/image19.png", Cm(10.5), Cm(7), width=Cm(8.9), height=Cm(6))
except Exception as e:
    pass

add_text(slide2, "PREVENCIÓN ACADÉMICA", Cm(10.5), Cm(14), Cm(8.9), Cm(1), font_size=Pt(14), bold=True, color=dark_blue)
add_text(slide2, "Gracias al motor de cálculos y reglas de evaluación (Ordenanza 04-2023), identificamos automáticamente a los estudiantes que necesitan apoyo y recuperación pedagógica temprana.", Cm(10.5), Cm(15.5), Cm(8.9), Cm(4), font_size=Pt(10))

# Panel 6 (Right)
shape = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(19.8), Cm(0), Cm(9.9), Cm(21))
shape.fill.solid()
shape.fill.fore_color.rgb = dark_blue
shape.line.fill.background()

add_text(slide2, "COMPONENTES DEL SISTEMA", Cm(20.5), Cm(2), Cm(8.5), Cm(1), font_size=Pt(16), bold=True, color=cyan)

add_text(slide2, "Motor de Cálculos Normativo\nEvaluación automática de competencias, completivas (50/50) y extraordinarias (30/70).", Cm(20.5), Cm(5), Cm(8.5), Cm(2), font_size=Pt(10), color=white)
add_text(slide2, "Generación de Boletines\nExportación masiva de boletines idénticos al formato oficial del MINERD.", Cm(20.5), Cm(8), Cm(8.5), Cm(2), font_size=Pt(10), color=white)
add_text(slide2, "Dashboard Directivo\nMétricas de rendimiento por sección, grado y docente, visibilidad total.", Cm(20.5), Cm(11), Cm(8.5), Cm(2), font_size=Pt(10), color=white)
add_text(slide2, "Alertas Tempranas\nDetección de estudiantes con ausentismo o con peligro de repitencia escolar.", Cm(20.5), Cm(14), Cm(8.5), Cm(2), font_size=Pt(10), color=white)

prs.save('Brochure_SIGERA.pptx')
